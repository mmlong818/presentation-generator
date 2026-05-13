import io
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Any

from models import Deck
from theme import ThemeColors, get_theme

logger = logging.getLogger(__name__)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def add_bg(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


def add_rect(slide: Any, x: float, y: float, w: float, h: float,
             fill_hex: str | None = None,
             line_hex: str | None = None,
             line_width_pt: float = 0) -> Any:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    else:
        shape.fill.background()
    if line_hex and line_width_pt > 0:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide: Any, x: float, y: float, w: float, h: float,
                text: str, font_name: str, font_size_pt: float,
                color_hex: str, bold: bool = False, italic: bool = False,
                align: str = "left", wrap: bool = True) -> Any:
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = hex_to_rgb(color_hex)
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_line(slide: Any, x1: float, y1: float, x2: float, y2: float,
             color_hex: str, width_pt: float = 1.0) -> Any:
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = hex_to_rgb(color_hex)
    connector.line.width = Pt(width_pt)
    return connector


def add_eyebrow(slide: Any, text: str, t: ThemeColors,
                x: float = 0.83, y: float = 0.5, w: float = 11.0) -> None:
    if not text:
        return
    add_textbox(slide, x, y, w, 0.4, text.upper(),
                t.font_body, 14, t.muted, bold=False)


def render_slide(slide: Any, data: dict, t: ThemeColors) -> None:
    layout_type = data.get("type", "diagram")
    from renderers import RENDERERS
    renderer = RENDERERS.get(layout_type)
    if renderer:
        renderer(slide, data, t)
    else:
        add_bg(slide, t.bg)
        add_textbox(slide, 0.83, 2.5, 11.0, 2.0,
                    f"[{layout_type}] — layout not rendered",
                    t.font_body, 28, t.muted, align="center")


def build_pptx(deck: Deck) -> bytes:
    t = get_theme(deck.theme)
    if deck.brand and deck.brand.accent:
        t.accent = deck.brand.accent

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    for slide_data in deck.slides:
        slide = prs.slides.add_slide(blank_layout)
        render_slide(slide, slide_data, t)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
